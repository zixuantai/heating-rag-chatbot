import os
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation


def parse_document(file_content, file_name):
    """根据文件类型解析文档内容"""
    ext = os.path.splitext(file_name)[1].lower()
    
    if ext == '.pdf':
        return parse_pdf(file_content)
    elif ext == '.docx':
        return parse_docx(file_content)
    elif ext == '.pptx':
        return parse_pptx(file_content)
    elif ext == '.txt':
        return file_content.decode('utf-8', errors='ignore')
    else:
        raise ValueError(f"不支持的文件类型: {ext}")


def parse_pdf(file_content):
    """解析PDF文档"""
    with open('temp.pdf', 'wb') as f:
        f.write(file_content)
    
    text = ""
    try:
        with open('temp.pdf', 'rb') as f:
            reader = PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception as e:
        text = f"[PDF解析警告: {str(e)}]"
    finally:
        if os.path.exists('temp.pdf'):
            os.remove('temp.pdf')
    
    return text


def parse_docx(file_content):
    """解析Word文档"""
    with open('temp.docx', 'wb') as f:
        f.write(file_content)
    
    text = ""
    try:
        doc = DocxDocument('temp.docx')
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
    except Exception as e:
        text = f"[Word解析警告: {str(e)}]"
    finally:
        if os.path.exists('temp.docx'):
            os.remove('temp.docx')
    
    return text


def parse_pptx(file_content):
    """解析PPT文档"""
    with open('temp.pptx', 'wb') as f:
        f.write(file_content)
    
    text = ""
    try:
        prs = Presentation('temp.pptx')
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + "\n"
    except Exception as e:
        text = f"[PPT解析警告: {str(e)}]"
    finally:
        if os.path.exists('temp.pptx'):
            os.remove('temp.pptx')
    
    return text
